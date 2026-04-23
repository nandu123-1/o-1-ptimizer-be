"""Generate O(1)ptimizer.pptx — a 7-slide pitch deck.

Usage:
    pip install python-pptx
    python scripts/make_ppt.py
    # Output: O(1)ptimizer.pptx in the repo root.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Dark theme palette
BG = RGBColor(0x0B, 0x12, 0x1C)
PANEL = RGBColor(0x14, 0x20, 0x2E)
ACCENT = RGBColor(0x2D, 0xD4, 0xBF)  # teal
ACCENT_DIM = RGBColor(0x14, 0xB8, 0xA6)
INK = RGBColor(0xE6, 0xF4, 0xF1)
INK_MUTED = RGBColor(0x94, 0xA3, 0xB8)
DANGER = RGBColor(0xF4, 0x3F, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, left, top, width, height, text, *, font_size=18,
             bold=False, color=INK, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_panel(slide, left, top, width, height, *, fill=PANEL, border=ACCENT_DIM):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_pill(slide, left, top, text, *, bg=ACCENT, fg=BG):
    width = Inches(2.1)
    height = Inches(0.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Segoe UI"
    return shape


def add_eyebrow(slide, top, text):
    return add_text(slide, Inches(0.7), top, Inches(10), Inches(0.35),
                    text.upper(), font_size=12, bold=True, color=ACCENT,
                    font_name="Consolas")


def add_title(slide, top, text, *, size=40):
    return add_text(slide, Inches(0.7), top, Inches(12), Inches(1.2),
                    text, font_size=size, bold=True, color=INK)


def add_body_bullets(slide, left, top, width, height, bullets,
                     *, font_size=16, gap=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)

    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = gap
        p.space_after = Pt(6)
        dot = p.add_run()
        dot.text = "  "
        dot_bullet = p.add_run()
        dot_bullet.text = "■  "
        dot_bullet.font.size = Pt(font_size)
        dot_bullet.font.color.rgb = ACCENT
        dot_bullet.font.name = "Segoe UI"
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.color.rgb = INK
        r.font.name = "Segoe UI"
    return tb


# ---------- Slides ----------

def slide1_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.7), Inches(1.0),
                                 Inches(0.08), Inches(2.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text(slide, Inches(0.95), Inches(1.0), Inches(10), Inches(0.5),
             "PROTOTYPE  ·  APRIL 2026", font_size=13, bold=True,
             color=ACCENT, font_name="Consolas")

    add_text(slide, Inches(0.95), Inches(1.55), Inches(12), Inches(1.6),
             "O(1)ptimizer", font_size=72, bold=True, color=INK)

    add_text(slide, Inches(0.95), Inches(3.2), Inches(12), Inches(0.8),
             "AI-powered C++ complexity reduction —",
             font_size=24, color=INK)
    add_text(slide, Inches(0.95), Inches(3.75), Inches(12), Inches(0.8),
             "from brute force to optimal in one click.",
             font_size=24, color=INK_MUTED)

    # Stack pills
    pills = [
        ("Next.js 16", ACCENT),
        ("FastAPI", ACCENT),
        ("Gemini", ACCENT),
        ("CrewAI", ACCENT),
    ]
    x = Inches(0.95)
    for label, color in pills:
        add_pill(slide, x, Inches(5.2), label, bg=color, fg=BG)
        x += Inches(2.25)

    add_text(slide, Inches(0.95), Inches(6.7), Inches(12), Inches(0.4),
             "o1ptimizer.netlify.app", font_size=14, color=INK_MUTED,
             font_name="Consolas")


def slide2_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_eyebrow(slide, Inches(0.8), "THE PROBLEM")
    add_title(slide, Inches(1.2), "Most student C++ ships at O(n²) or worse.")

    bullets = [
        "Nested loops when a hash map would be O(n).",
        "Naive data structures — linear scans where a set would do.",
        "Re-computed work inside loops that could be hoisted.",
        "IDEs catch syntax, not asymptotic cost. There is no tight feedback loop",
        "between ‘does it compile’ and ‘does it scale’.",
    ]
    add_body_bullets(slide, Inches(0.7), Inches(2.9), Inches(12), Inches(4),
                     bullets, font_size=20)


def slide3_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_eyebrow(slide, Inches(0.8), "THE SOLUTION")
    add_title(slide, Inches(1.2),
              "Paste code → get an optimized rewrite + reasons why.")

    # Three cards
    card_w = Inches(3.9)
    card_h = Inches(3.2)
    top = Inches(3.2)
    gap = Inches(0.25)
    left0 = Inches(0.7)

    cards = [
        ("1. Submit",
         "Paste, upload a .cpp, or snap a photo of handwritten code. Gemini vision OCRs the image."),
        ("2. Optimize",
         "Single-shot Gemini call returns strict JSON with optimized C++, Big-O before/after, and algorithm picks."),
        ("3. Visualize",
         "Recharts line chart plots brute-force vs optimized operations across input sizes 10 → 100k."),
    ]
    for i, (head, body) in enumerate(cards):
        left = left0 + (card_w + gap) * i
        add_panel(slide, left, top, card_w, card_h)
        add_text(slide, left + Inches(0.3), top + Inches(0.25),
                 card_w - Inches(0.6), Inches(0.5),
                 head, font_size=20, bold=True, color=ACCENT)
        add_text(slide, left + Inches(0.3), top + Inches(0.9),
                 card_w - Inches(0.6), card_h - Inches(1.2),
                 body, font_size=15, color=INK)


def slide4_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_eyebrow(slide, Inches(0.8), "ARCHITECTURE")
    add_title(slide, Inches(1.2), "Two-service split, one prompt.")

    # Flow diagram: Browser -> Netlify (frontend) -> Render (backend) -> Gemini
    blocks = [
        ("Browser", "Next.js 16 + Monaco\nRecharts"),
        ("Netlify", "Static build\nCDN edge"),
        ("Render", "FastAPI\nfailover engine"),
        ("Gemini", "flash family\n2 keys × 6 models"),
    ]
    block_w = Inches(2.6)
    block_h = Inches(2.0)
    top = Inches(3.2)
    gap = Inches(0.45)
    total_w = block_w * 4 + gap * 3
    start_left = (SLIDE_W - total_w) / 2

    for i, (head, body) in enumerate(blocks):
        left = start_left + (block_w + gap) * i
        add_panel(slide, left, top, block_w, block_h)
        add_text(slide, left + Inches(0.25), top + Inches(0.25),
                 block_w - Inches(0.5), Inches(0.5),
                 head, font_size=18, bold=True, color=ACCENT)
        add_text(slide, left + Inches(0.25), top + Inches(0.85),
                 block_w - Inches(0.5), block_h - Inches(1.1),
                 body, font_size=13, color=INK_MUTED,
                 font_name="Consolas")

        if i < len(blocks) - 1:
            arrow_left = left + block_w + Inches(0.05)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left, top + block_h / 2 - Inches(0.15),
                Inches(0.35), Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    add_text(slide, Inches(0.7), Inches(5.8), Inches(12), Inches(0.6),
             "Browser calls Render directly via NEXT_PUBLIC_BACKEND_API_BASE_URL,",
             font_size=14, color=INK_MUTED)
    add_text(slide, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6),
             "sidestepping Netlify’s 10s function timeout during Render cold-starts.",
             font_size=14, color=INK_MUTED)


def slide5_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_eyebrow(slide, Inches(0.8), "LIVE DEMO")
    add_title(slide, Inches(1.2), "Two Sum: O(n²) → O(n) in 25 seconds.")

    # Two code panels side-by-side
    top = Inches(3.0)
    height = Inches(3.3)
    gap = Inches(0.3)
    col_w = (SLIDE_W - Inches(1.4) - gap) / 2

    left_panel_x = Inches(0.7)
    right_panel_x = left_panel_x + col_w + gap

    # Left
    add_panel(slide, left_panel_x, top, col_w, height, border=DANGER)
    add_text(slide, left_panel_x + Inches(0.3), top + Inches(0.2),
             col_w - Inches(0.6), Inches(0.4),
             "BEFORE  —  O(n²)", font_size=14, bold=True,
             color=DANGER, font_name="Consolas")
    before_code = (
        "for (int i = 0; i < n; ++i) {\n"
        "  for (int j = i+1; j < n; ++j) {\n"
        "    if (v[i] + v[j] == t) {\n"
        "      return {i, j};\n"
        "    }\n"
        "  }\n"
        "}"
    )
    add_text(slide, left_panel_x + Inches(0.3), top + Inches(0.75),
             col_w - Inches(0.6), height - Inches(0.9),
             before_code, font_size=13, color=INK, font_name="Consolas")

    # Right
    add_panel(slide, right_panel_x, top, col_w, height, border=ACCENT)
    add_text(slide, right_panel_x + Inches(0.3), top + Inches(0.2),
             col_w - Inches(0.6), Inches(0.4),
             "AFTER  —  O(n)", font_size=14, bold=True,
             color=ACCENT, font_name="Consolas")
    after_code = (
        "unordered_map<int,int> seen;\n"
        "for (int i = 0; i < n; ++i) {\n"
        "  int need = t - v[i];\n"
        "  auto it = seen.find(need);\n"
        "  if (it != seen.end())\n"
        "    return {it->second, i};\n"
        "  seen[v[i]] = i;\n"
        "}"
    )
    add_text(slide, right_panel_x + Inches(0.3), top + Inches(0.75),
             col_w - Inches(0.6), height - Inches(0.9),
             after_code, font_size=13, color=INK, font_name="Consolas")

    # Speedup callout
    add_panel(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.75))
    add_text(slide, Inches(0.95), Inches(6.65), Inches(12), Inches(0.5),
             "estimated_speedup_ratio ≈ 10,000× at n = 10,000",
             font_size=16, bold=True, color=ACCENT, font_name="Consolas")


def slide6_highlights(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_eyebrow(slide, Inches(0.8), "ENGINEERING HIGHLIGHTS")
    add_title(slide, Inches(1.2), "Production-ready edges.")

    bullets = [
        "Single-shot Gemini path cuts LLM calls ~15× vs multi-agent — fits free-tier quota.",
        "Automatic key + model rotation on 429 RESOURCE_EXHAUSTED (2 keys × 6 models = 12 combos).",
        "Strict Pydantic contract validates every response before it ever reaches the UI.",
        "Image-to-code: Gemini vision extracts C++ from a photo of handwritten pseudocode.",
        "Self-diagnosing frontend — non-JSON responses surface HTTP status + body preview.",
        "Crew fallback: any multi-agent failure falls through to single-shot so one bad call never kills a request.",
    ]
    add_body_bullets(slide, Inches(0.7), Inches(2.9), Inches(12), Inches(4.3),
                     bullets, font_size=17)


def slide7_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_eyebrow(slide, Inches(0.8), "ROADMAP")
    add_title(slide, Inches(1.2), "What’s next.")

    bullets = [
        "Execute optimized code in a sandboxed g++ container and diff outputs vs original.",
        "Multi-language support — Python, Java, Rust.",
        "Batch mode for graded assignments (CSV in, annotated CSV out).",
        "“Explain like I’m a first-year CS student” toggle for pedagogical use.",
        "Per-function complexity heat-map overlay on the Monaco editor.",
    ]
    add_body_bullets(slide, Inches(0.7), Inches(2.9), Inches(12), Inches(3.8),
                     bullets, font_size=17)

    add_panel(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.8))
    add_text(slide, Inches(0.95), Inches(6.55), Inches(12), Inches(0.5),
             "Ship code that runs in the complexity class it deserves.",
             font_size=18, bold=True, color=ACCENT)


# ---------- Build ----------

def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide1_title(prs)
    slide2_problem(prs)
    slide3_solution(prs)
    slide4_architecture(prs)
    slide5_demo(prs)
    slide6_highlights(prs)
    slide7_roadmap(prs)

    prs.save(output)
    print(f"Wrote {output}")


if __name__ == "__main__":
    root = Path(__file__).resolve().parent.parent
    build(root / "O(1)ptimizer.pptx")
